"""Styled deck builder — turns a slide outline (+ optional AI images) into a
designed .pptx using python-pptx.

The default template renders flat white bullet slides. This module instead
draws every slide on a blank layout with a small design system: a colour
palette, a clean font, accent bars, section dividers, two-column image
layouts, and cover-cropped images. The result is rasterised to PNGs by the
existing conversion pipeline, so the viewer + signed-URL model are unchanged.

Images are optional per slide: when an image is missing (generation failed or
disabled) the layout falls back to a full-width text panel, so a deck always
builds.
"""

from __future__ import annotations

import io

from django.conf import settings
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 16:9 canvas.
_W = Inches(13.333)
_H = Inches(7.5)


def _palette() -> dict[str, RGBColor]:
    """Brand palette, overridable via settings.DECK_PALETTE (hex strings)."""
    defaults = {
        "primary": "0F4C81",   # deep blue — headers, title bg
        "accent": "F6A609",    # amber — bars, bullet markers
        "dark": "0F172A",      # near-black — body text
        "muted": "64748B",     # slate — captions/footer
        "light": "F1F5F9",     # light slate — section/quote bg
        "white": "FFFFFF",
    }
    overrides = getattr(settings, "DECK_PALETTE", None) or {}
    merged = {**defaults, **overrides}
    return {k: RGBColor.from_string(v) for k, v in merged.items()}


FONT = getattr(settings, "DECK_FONT", "DejaVu Sans")


# --------------------------------------------------------------------------
# low-level drawing helpers
# --------------------------------------------------------------------------
def _rect(slide, x, y, w, h, color, *, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size,
    color,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def _bullets(slide, x, y, w, h, bullets, pal, *, size=18):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        marker = p.add_run()
        marker.text = "●  "  # ● accent dot
        marker.font.size = Pt(size - 4)
        marker.font.name = FONT
        marker.font.color.rgb = pal["accent"]
        run = p.add_run()
        run.text = str(bullet)
        run.font.size = Pt(size)
        run.font.name = FONT
        run.font.color.rgb = pal["dark"]
    return box


def _cover_image(slide, image_bytes, x, y, w, h):
    """Place an image filling the box, cover-cropped (no distortion)."""
    from PIL import Image

    iw, ih = Image.open(io.BytesIO(image_bytes)).size
    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), x, y, width=w, height=h)
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:  # image too wide -> trim sides
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif img_ratio < box_ratio:  # image too tall -> trim top/bottom
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def _footer(slide, index, total, deck_title, pal):
    _text(
        slide,
        Inches(0.9),
        Inches(7.02),
        Inches(9),
        Inches(0.35),
        deck_title,
        size=10,
        color=pal["muted"],
    )
    _text(
        slide,
        Inches(11.6),
        Inches(7.02),
        Inches(1.2),
        Inches(0.35),
        f"{index} / {total}",
        size=10,
        color=pal["muted"],
        align=PP_ALIGN.RIGHT,
    )


# --------------------------------------------------------------------------
# per-slide layouts
# --------------------------------------------------------------------------
def _title_slide(slide, title, subtitle, image, pal):
    _rect(slide, 0, 0, _W, _H, pal["primary"])
    if image:
        # framed image panel on the right
        _rect(slide, Inches(8.75), Inches(0.95), Inches(4.0), Inches(5.6), pal["accent"])
        _cover_image(slide, image, Inches(8.9), Inches(1.1), Inches(3.7), Inches(5.3))
    _rect(slide, Inches(0.9), Inches(2.5), Inches(1.6), Inches(0.16), pal["accent"])
    _text(
        slide, Inches(0.9), Inches(2.8), Inches(7.4), Inches(2.6),
        title, size=44, color=pal["white"], bold=True,
    )
    if subtitle:
        _text(
            slide, Inches(0.9), Inches(4.9), Inches(7.4), Inches(1.2),
            subtitle, size=20, color=RGBColor.from_string("C7D2E0"),
        )


def _section_slide(slide, title, image, pal):
    _rect(slide, 0, 0, _W, _H, pal["light"])
    _rect(slide, 0, 0, Inches(0.45), _H, pal["accent"])
    _text(
        slide, Inches(1.3), 0, Inches(10.5), _H,
        title, size=40, color=pal["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )


def _quote_slide(slide, title, bullets, image, pal):
    _rect(slide, 0, 0, _W, _H, pal["light"])
    _text(
        slide, Inches(0.8), Inches(0.9), Inches(3), Inches(2),
        "“", size=140, color=pal["accent"], bold=True,
    )
    _text(
        slide, Inches(1.6), Inches(2.4), Inches(10.1), Inches(3),
        title, size=30, color=pal["dark"], italic=True, anchor=MSO_ANCHOR.MIDDLE,
    )
    if bullets:
        _text(
            slide, Inches(1.6), Inches(5.4), Inches(10.1), Inches(0.8),
            f"— {bullets[0]}", size=18, color=pal["primary"], bold=True,
        )


def _content_slide(slide, title, bullets, image, pal):
    _rect(slide, 0, 0, _W, _H, pal["white"])
    # header: accent bar + title + divider
    _rect(slide, Inches(0.9), Inches(0.7), Inches(1.1), Inches(0.14), pal["accent"])
    _text(
        slide, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0),
        title, size=30, color=pal["primary"], bold=True,
    )
    _rect(slide, Inches(0.9), Inches(1.95), Inches(11.5), Pt(1.5), pal["light"])

    if image:
        _bullets(slide, Inches(0.9), Inches(2.3), Inches(6.5), Inches(4.3), bullets, pal)
        # framed image panel on the right
        _rect(slide, Inches(7.75), Inches(2.2), Inches(4.6), Inches(4.35), pal["light"])
        _cover_image(slide, image, Inches(7.9), Inches(2.35), Inches(4.3), Inches(4.05))
    else:
        _bullets(
            slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.3),
            bullets, pal, size=20,
        )


# --------------------------------------------------------------------------
# public API
# --------------------------------------------------------------------------
def build_pptx(outline: dict, images: list[bytes | None]) -> bytes:
    """Render a validated outline (+ per-slide images) into styled .pptx bytes.

    `images[i]` aligns with `outline["slides"][i]`; any entry may be None.
    """
    pal = _palette()
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    blank = prs.slide_layouts[6]

    slides = outline.get("slides") or []
    total = len(slides)
    deck_title = outline.get("title") or "Taqdimot"

    for idx, entry in enumerate(slides):
        image = images[idx] if idx < len(images) else None
        layout = (entry.get("layout") or "content").lower()
        title = str(entry.get("title", ""))
        bullets = [str(b) for b in (entry.get("bullets") or [])]
        slide = prs.slides.add_slide(blank)

        if layout == "title" or idx == 0:
            _title_slide(slide, title or deck_title, outline.get("subtitle", ""), image, pal)
        elif layout == "section":
            _section_slide(slide, title, image, pal)
        elif layout == "quote":
            _quote_slide(slide, title, bullets, image, pal)
        else:
            _content_slide(slide, title, bullets, image, pal)
            _footer(slide, idx + 1, total, deck_title, pal)

        notes = entry.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
